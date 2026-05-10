"""
Script to generate PowerPoint presentation from ITMP_slides.md
Each "# Dia X:" section creates one slide with:
  - Title from "Címsor"
  - Content from "Szöveges tartalom"
  - Speaker notes from "Megjegyzések az előadónak"
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def parse_markdown_file(filename):
    """Parse the markdown file and extract sections."""
    with open(filename, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split by main sections (# Dia X:)
    sections = re.split(r'\n(?=# Dia \d+:)', content)
    sections = [s.strip() for s in sections if s.strip() and s.startswith('# Dia')]
    
    parsed_sections = []
    for section in sections:
        lines = section.split('\n')
        
        # Extract the main section title
        main_title = lines[0].replace('# ', '').strip()
        
        # Initialize section data
        section_data = {
            'main_title': main_title,
            'címsor': '',
            'szoveges_tartalom': '',
            'illusztracio': '',
            'megjegyzes': ''
        }
        
        current_subsection = None
        current_content = []
        
        for line in lines[1:]:
            if line.startswith('### '):
                # Save previous subsection
                if current_subsection:
                    section_data[current_subsection] = '\n'.join(current_content).strip()
                
                # Determine new subsection
                header = line.replace('### ', '').strip()
                if 'Címsor' in header:
                    current_subsection = 'címsor'
                elif 'Szöveges tartalom' in header:
                    current_subsection = 'szoveges_tartalom'
                elif 'Illusztráció' in header:
                    current_subsection = 'illusztracio'
                elif 'Megjegyzés' in header:
                    current_subsection = 'megjegyzes'
                else:
                    current_subsection = None
                
                current_content = []
            else:
                if current_subsection:
                    current_content.append(line)
        
        # Save last subsection
        if current_subsection:
            section_data[current_subsection] = '\n'.join(current_content).strip()
        
        parsed_sections.append(section_data)
    
    return parsed_sections

def add_formatted_text(text_frame, text, base_font_size=14):
    """Add text with markdown-style formatting."""
    text_frame.clear()
    
    paragraphs = text.split('\n')
    for i, para_text in enumerate(paragraphs):
        if not para_text.strip():
            continue
            
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        # Handle bold text (**text**)
        if '**' in para_text:
            parts = re.split(r'(\*\*.*?\*\*)', para_text)
            for part in parts:
                if part.startswith('**') and part.endswith('**'):
                    run = p.add_run()
                    run.text = part[2:-2]
                    run.font.bold = True
                    run.font.size = Pt(base_font_size)
                    run.font.name = 'Calibri'
                else:
                    run = p.add_run()
                    run.text = part
                    run.font.size = Pt(base_font_size)
                    run.font.name = 'Calibri'
        else:
            p.text = para_text
            p.font.size = Pt(base_font_size)
            p.font.name = 'Calibri'
        
        # Bullets for lines starting with •
        if para_text.strip().startswith('•'):
            p.level = 0
            p.text = para_text.strip()[1:].strip()
        elif para_text.strip().startswith('✓'):
            p.level = 0
            p.text = para_text.strip()[1:].strip()
        
        # Indentation for continuation lines
        if para_text.startswith('  ') and not para_text.strip().startswith('•'):
            p.level = 1

def create_presentation(sections, output_filename):
    """Create PowerPoint presentation from parsed sections."""
    prs = Presentation()
    
    # Set slide dimensions (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Process each section as one slide
    for section in sections:
        # Create slide with title and content layout
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Title box
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.text = section['címsor']
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.name = 'Calibri'
        title_p.alignment = PP_ALIGN.CENTER
        
        # Content box
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9), Inches(5.4))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        # Add formatted content
        content_text = section['szoveges_tartalom']
        add_formatted_text(content_frame, content_text, base_font_size=16)
        
        # Add speaker notes
        if section['megjegyzes']:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = section['megjegyzes']
    
    # Save presentation
    prs.save(output_filename)
    print(f"✓ Prezentáció sikeresen létrehozva: {output_filename}")
    print(f"✓ Összes slide: {len(prs.slides)}")

if __name__ == '__main__':
    input_file = 'ITMP_slides.md'
    output_file = 'ITMP_slides.pptx'
    
    print("Markdown fájl feldolgozása...")
    sections = parse_markdown_file(input_file)
    print(f"✓ {len(sections)} dia beolvasva")
    
    print("PowerPoint prezentáció generálása...")
    create_presentation(sections, output_file)
