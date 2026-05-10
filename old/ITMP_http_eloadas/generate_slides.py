"""
Script to generate PowerPoint presentation from ITMP_task_slides.md
Each "Feladat X:" section creates 3 slides:
  - Slide 1: Címsor + Feladat leírása + Tanulság, mire jó
  - Slide 2: Példa prompt 1
  - Slide 3: Példa prompt 2
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def parse_markdown_file(filename):
    """Parse the markdown file and extract sections."""
    with open(filename, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split by main sections (# Feladat X:)
    sections = re.split(r'\n(?=# Feladat \d+:)', content)
    sections = [s.strip() for s in sections if s.strip() and s.startswith('# Feladat')]
    
    parsed_sections = []
    for section in sections:
        lines = section.split('\n')
        
        # Extract the main title
        main_title = lines[0].replace('# ', '').strip()
        
        # Initialize section data
        section_data = {
            'main_title': main_title,
            'címsor': '',
            'feladat_leiras': '',
            'tanulsag': '',
            'megjegyzes': '',
            'pelda_1': '',
            'pelda_2': ''
        }
        
        current_subsection = None
        current_content = []
        
        for line in lines[1:]:
            if line.startswith('### '):
                # Save previous subsection
                if current_subsection:
                    section_data[current_subsection] = '\n'.join(current_content).strip()
                
                # Determine new subsection
                header = line.replace('### ', '').strip()
                if 'Címsor' in header:
                    current_subsection = 'címsor'
                elif 'Feladat leírása' in header:
                    current_subsection = 'feladat_leiras'
                elif 'Tanulság' in header:
                    current_subsection = 'tanulsag'
                elif 'Megjegyzés' in header:
                    current_subsection = 'megjegyzes'
                elif 'Példa prompt 1' in header:
                    current_subsection = 'pelda_1'
                elif 'Példa prompt 2' in header:
                    current_subsection = 'pelda_2'
                else:
                    current_subsection = None
                
                current_content = []
            else:
                if current_subsection:
                    current_content.append(line)
        
        # Save last subsection
        if current_subsection:
            section_data[current_subsection] = '\n'.join(current_content).strip()
        
        parsed_sections.append(section_data)
    
    return parsed_sections

def add_text_to_textbox(textbox, text, font_size=14, bold_title=False):
    """Add text to a textbox with formatting."""
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.clear()
    
    paragraphs = text.split('\n')
    for i, para_text in enumerate(paragraphs):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        p.text = para_text
        p.font.size = Pt(font_size)
        p.font.name = 'Calibri'
        
        # Bold formatting for specific patterns
        if bold_title and (para_text.startswith('**') or para_text.startswith('•')):
            p.font.bold = True

def create_presentation(sections, output_filename):
    """Create PowerPoint presentation from parsed sections."""
    prs = Presentation()
    
    # Set slide dimensions (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = "ITMP HTTP Előadás"
    subtitle.text = "AI Promptolási Feladatok"
    
    # Process each section
    for section in sections:
        # Slide 1: Címsor + Feladat leírása + Tanulság
        slide1_layout = prs.slide_layouts[5]  # Blank layout
        slide1 = prs.slides.add_slide(slide1_layout)
        
        # Title
        title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = section['címsor']
        title_p.font.size = Pt(28)
        title_p.font.bold = True
        title_p.font.name = 'Calibri'
        title_p.alignment = PP_ALIGN.CENTER
        
        # Feladat leírása
        desc_box = slide1.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(3))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        desc_p = desc_frame.paragraphs[0]
        desc_p.text = section['feladat_leiras']
        desc_p.font.size = Pt(12)
        desc_p.font.name = 'Calibri'
        
        # Tanulság
        tanu_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2.5))
        tanu_frame = tanu_box.text_frame
        tanu_frame.word_wrap = True
        tanu_p = tanu_frame.paragraphs[0]
        tanu_p.text = section['tanulsag']
        tanu_p.font.size = Pt(12)
        tanu_p.font.name = 'Calibri'
        
        # Slide 2: Példa prompt 1
        slide2_layout = prs.slide_layouts[5]
        slide2 = prs.slides.add_slide(slide2_layout)
        
        # Title for slide 2
        title2_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title2_frame = title2_box.text_frame
        title2_p = title2_frame.paragraphs[0]
        title2_p.text = f"{section['címsor']} - Példa 1"
        title2_p.font.size = Pt(24)
        title2_p.font.bold = True
        title2_p.font.name = 'Calibri'
        
        # Példa prompt 1 content
        pelda1_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(6))
        pelda1_frame = pelda1_box.text_frame
        pelda1_frame.word_wrap = True
        pelda1_p = pelda1_frame.paragraphs[0]
        pelda1_p.text = section['pelda_1']
        pelda1_p.font.size = Pt(10)
        pelda1_p.font.name = 'Calibri'
        
        # Slide 3: Példa prompt 2
        slide3_layout = prs.slide_layouts[5]
        slide3 = prs.slides.add_slide(slide3_layout)
        
        # Title for slide 3
        title3_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title3_frame = title3_box.text_frame
        title3_p = title3_frame.paragraphs[0]
        title3_p.text = f"{section['címsor']} - Példa 2"
        title3_p.font.size = Pt(24)
        title3_p.font.bold = True
        title3_p.font.name = 'Calibri'
        
        # Példa prompt 2 content
        pelda2_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(6))
        pelda2_frame = pelda2_box.text_frame
        pelda2_frame.word_wrap = True
        pelda2_p = pelda2_frame.paragraphs[0]
        pelda2_p.text = section['pelda_2']
        pelda2_p.font.size = Pt(10)
        pelda2_p.font.name = 'Calibri'
    
    # Save presentation
    prs.save(output_filename)
    print(f"✓ Prezentáció sikeresen létrehozva: {output_filename}")
    print(f"✓ Összes slide: {len(prs.slides)}")

if __name__ == '__main__':
    input_file = 'ITMP_task_slides.md'
    output_file = 'ITMP_task_slides.pptx'
    
    print("Markdown fájl feldolgozása...")
    sections = parse_markdown_file(input_file)
    print(f"✓ {len(sections)} szakasz beolvasva")
    
    print("PowerPoint prezentáció generálása...")
    create_presentation(sections, output_file)
